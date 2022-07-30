#for creating ppt.

from pptx import Presentation
import pptx
import sys
# For watermark Image.
from wand.image import Image

# Imprting Inches for adjesting image.
from pptx.util import Inches

#importing os for intrecting with oprating system.
import os


# Creat ppt.
prs = Presentation()


class MySlide:

    def __init__(self,data):

        # des = 'C:/Users/Lenovo/Desktop/Test'
        des = str(sys.argv[1])
        # Checking jpg format of files in directory.

        for w in os.listdir(des):

            if w.endswith('.jpg'):
                
                with Image(filename = w) as image:
            
                    # Import the watermark image
                    with Image(filename ='nike_black.png') as water:
                
                        water.resize(width=1000, height=600,)
                         # Clone the image in order to process
                        with image.clone() as watermark:

                            # Invoke watermark function with watermark image, transparency as 0.5s
                            # left as 10 and top as 20
                            watermark.watermark(water, 0,5,5,)
                            # Save the image
                            p = str(w).rpartition(".")
                            fname= str(p[0])+"watermark.jpg"
                            watermark.save(filename = fname)         

                    # Adding slides with image.

                    self.blank_slide_layout = prs.slide_layouts[data[2]]
                    self.slide = prs.slides.add_slide(self.blank_slide_layout)

                    self.title = self.slide.shapes.title
                    self.title.text = data[0]
                    self.subtitle = self.slide.placeholders[1]
                    self.subtitle.text = data[1]

                    # Addjesting image.
                
                    left = top = Inches (3)

                    left = Inches(1)
                    height = Inches(4)

                    pic = self.slide.shapes.add_picture(fname,left, top, height = height)

# Adding Title and subtitle in each slides
slides = [
            
["Title of the Presentation",
"Subtitle",
1],

]

for es in slides:
    MySlide(es)
# Saveing ppt.
prs.save("test.pptx")
# Open ppt .
os.startfile("test.pptx")





